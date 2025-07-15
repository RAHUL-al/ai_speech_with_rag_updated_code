import ujson as json
import redis
from datetime import datetime, timedelta
from fastapi import FastAPI, Depends, HTTPException, status, Request, WebSocket, WebSocketDisconnect, UploadFile, File, Form, Query
from fastapi.middleware.cors import CORSMiddleware
from pydub import AudioSegment
from dotenv import load_dotenv
from urllib.parse import parse_qs, unquote
from concurrent.futures import ThreadPoolExecutor
from fastapi.responses import FileResponse, JSONResponse
import asyncio
import time
from pdf2image import convert_from_path
from typing import List
import os, shutil, uuid
from PIL import Image
import google.generativeai as genai
from pinecone import Pinecone
from auth import hash_password, verify_password, create_access_token
from langchain_huggingface import HuggingFaceEmbeddings
from langchain_community.vectorstores import Pinecone as PineconeVectorStore
import torch
import logging
from urllib.parse import parse_qs,unquote
from jose import JWTError, jwt
from ai_speech_module import Topic
from schemas import UserCreate, UserOut, UserUpdate, Token, LoginRequest, ForgotPasswordRequest, GeminiRequest, ChatRequest
from langchain_ollama import OllamaLLM
import re
from typing import List
from fastapi import FastAPI, UploadFile, File, Form
from pdf2image import convert_from_path
from PIL import Image
from docx2pdf import convert as docx_to_pdf
import os, uuid, shutil, logging, tempfile
from langchain.text_splitter import RecursiveCharacterTextSplitter
from fastapi import HTTPException, status
from langchain.prompts import PromptTemplate
from langchain.chains import RetrievalQA
from langchain_community.vectorstores import Pinecone as PineconeVectorStore
# from langchain_community.vectorstores import Pinecone
from PIL import Image, ImageDraw, ImageFont
from PIL import Image, ImageDraw, ImageFont
from pdf2image import convert_from_path
from fastapi import UploadFile, File, Form, HTTPException
import pandas as pd
import docx2txt
from pptx import Presentation
import requests
from bs4 import BeautifulSoup
from firestore_models import FirestoreEssay
from firebase import db

scraping_api_key = os.getenv("SCRAPINGDOG_API_KEY")

PINECONE_API_KEY = os.getenv("PINECONE_API_KEY")
PINECONE_ENV = os.getenv("PINECONE_ENVIRONMENT")
PINECONE_INDEX_NAME = os.getenv("PINECONE_INDEX_NAME")
GOOGLE_API_KEY = os.getenv("GOOGLE_API_KEY")
EMBEDDING_MODEL_NAME = "embaas/sentence-transformers-e5-large-v2"

SECRET_KEY = "jwt_secret_key"
ALGORITHM="HS256"


logging.basicConfig(
    filename='rag_log.log',
    level=logging.INFO,
    format='%(asctime)s - %(levelname)s - %(message)s'
)

load_dotenv()
app = FastAPI(title="FastAPI‑Firebase")

redis_client = redis.StrictRedis(
    host=os.getenv("REDIS_HOST"),
    port=int(os.getenv("REDIS_PORT")),
    username=os.getenv("REDIS_USERNAME"),
    password=os.getenv("REDIS_PASSWORD"),
    decode_responses=True
)

origins = ["http://localhost:3000","http://localhost:5173"]

app.add_middleware(CORSMiddleware, allow_origins=origins, allow_credentials=True, allow_methods=["*"], allow_headers=["*"])

def get_user_from_redis_session(request: Request):
    token = request.headers.get("Authorization")
    if not token or not token.startswith("Bearer "):
        raise HTTPException(status_code=401, detail="Missing or invalid token")
    token = token.split(" ")[1]
    session_data = redis_client.get(f"session:{token}")
    if not session_data:
        raise HTTPException(status_code=401, detail="Session expired or invalid")
    return json.loads(session_data)

@app.post("/register", response_model=UserOut)
def register(user: UserCreate):
    user_ref = db.collection("users").where("username", "==", user.username).stream()
    if any(user_ref):
        raise HTTPException(400, "Username or email already exists")

    doc_ref = db.collection("users").add({
        "username": user.username,
        "email": user.email,
        "password": hash_password(user.password)
    })
    user_id = doc_ref[1].id
    return UserOut(id=user_id, username=user.username, email=user.email)

@app.post("/login", response_model=Token)
def login(data: LoginRequest):
    docs = db.collection("users").where("username", "==", data.username).stream()
    user_doc = next(docs, None)
    if not user_doc or not verify_password(data.password, user_doc.to_dict()["password"]):
        raise HTTPException(status.HTTP_401_UNAUTHORIZED, detail="Bad credentials")

    token = create_access_token({"sub": user_doc.id})
    redis_client.setex(f"session:{token}", timedelta(hours=1), json.dumps({"user_id": user_doc.id, "username": data.username}))
    return Token(access_token=token, username=data.username)

@app.get("/logout")
def logout(request: Request):
    token = request.headers.get("Authorization")
    if token and token.startswith("Bearer "):
        redis_client.delete(f"session:{token.split(' ')[1]}")
    return {"detail": "Logged out"}

@app.get("/me", response_model=UserOut)
def me(user=Depends(get_user_from_redis_session)):
    doc = db.collection("users").document(user["user_id"]).get()
    if not doc.exists:
        raise HTTPException(404, "User not found")
    data = doc.to_dict()
    return UserOut(id=doc.id, username=data["username"], email=data["email"])



@app.post("/generate-prompt")
def generate_prompt(data: GeminiRequest, user=Depends(get_user_from_redis_session)):
    url = f"https://en.wikipedia.org/wiki/{data.topic}"
    api_endpoint = f"https://api.scrapingdog.com/scrape?api_key={scraping_api_key}&url={url}"
    response = requests.get(api_endpoint)

    if response.status_code == 200:
        soup = BeautifulSoup(response.text, "html.parser")
        for script in soup(["script", "style", "noscript"]):
            script.decompose()
        text = soup.get_text(separator="\n", strip=True)
    else:
        logging.info(f"Error: {response.status_code} - {response.text}")
        text = ""

    prompt = f"Generate a essay for a student in class {data.student_class} with a {data.accent} accent, on the topic '{data.topic}', and the mood is '{data.mood}' and give me essay should be less than 400 words and in response did not want \n\n or \n and also not want word count thanks you this type of stuff and used {text} content for as updated data from internet and which is helpful in created essay and please give me output in paragraph format only not in points."

    username = user.get("username")
    topic = Topic()
    response_text = topic.topic_data_model_for_Qwen(username, prompt)

    essay_data = FirestoreEssay(
        username =username,
        user_id=user["user_id"],
        student_class=data.student_class,
        accent=data.accent,
        topic=data.topic,
        mood=data.mood,
        content=response_text
    )

    write_time, doc_ref = db.collection("essays").add(essay_data.to_dict())
    essay_id = doc_ref.id 

    return {
        "response": response_text,
        "essay_id": essay_id
    }


@app.get("/overall-scoring-by-id")
async def overall_scoring_by_id(essay_id: str):
    topic = Topic()
    result = await topic.overall_scoring_by_id(essay_id)
    return result


TEMP_DIR = os.path.abspath("audio_folder")
os.makedirs(TEMP_DIR, exist_ok=True)

@app.websocket("/ws/audio")
async def audio_ws(websocket: WebSocket):
    await websocket.accept()
    query_params = parse_qs(websocket.url.query)
    logging.info(query_params, "=-=-=-=-=-=-=-=-=-=-=-=-=-=-=-=-=-=-=-")
    username = query_params.get("username", [None])[0]
    token = query_params.get("token", [None])[0]
    logging.info(username, "=---------------")
    logging.info("Raw token received:", token)

    if not username or not token:
        await websocket.close(code=4001)
        logging.info("Username or token missing.")
        return

    try:
        token = unquote(token)
        if token.startswith("Bearer "):
            token = token[len("Bearer "):]
        elif token.startswith("Bearer+"):
            token = token[len("Bearer+"):]

        logging.info("Cleaned JWT token:", token)

        payload = jwt.decode(token, SECRET_KEY, algorithms=[ALGORITHM])
        token_username = payload.get("sub")

        logging.info("Decoded username from token:", token_username)


    except JWTError as e:
        await websocket.close(code=4003)
        logging.info(f"[JWT Error] {e}")
        return


    logging.info(f"[WS] Authenticated connection from {username}")
    chunk_index = 0
    chunk_files = []
    text_output = []
    topic = Topic()

    date_str = datetime.now().strftime("%Y-%m-%d")
    user_dir = os.path.join(TEMP_DIR, username, date_str)
    os.makedirs(user_dir, exist_ok=True)

    final_output = os.path.join(user_dir, f"{username}_output.wav")
    transcript_path = os.path.join(user_dir, f"{username}_transcript.txt")

    if os.path.exists(final_output):
        os.remove(final_output)
    if os.path.exists(transcript_path):
        os.remove(transcript_path)

    loop = asyncio.get_event_loop()

    try:
        while True:
            message = await websocket.receive()

            if message["type"] == "websocket.disconnect":
                print(f"[WS] {username} disconnected.")
                break

            if message["type"] == "websocket.receive" and "bytes" in message:
                chunk_filename = os.path.join(user_dir, f"chunk_{chunk_index}.wav")
                audio = AudioSegment(
                    data=message["bytes"],
                    sample_width=2,
                    frame_rate=16000,
                    channels=1
                )
                audio.export(chunk_filename, format="wav")
                chunk_files.append(chunk_filename)
                print("print chunk files",chunk_files)

                results = await asyncio.gather(
                    asyncio.to_thread(topic.speech_to_text, chunk_filename),
                    asyncio.to_thread(topic.detect_emotion, chunk_filename),
                    asyncio.to_thread(topic.fluency_scoring, chunk_filename),
                    asyncio.to_thread(topic.pronunciation_scoring, chunk_filename),
                    asyncio.to_thread(topic.silvero_vad, chunk_filename),
                )
                logging.info("results ---------->",results)
                transcribed_text, emotion, fluency, pronunciation, silvero = results

                text_output.append(transcribed_text)
                print(f"[Chunk {chunk_index}] Transcribed: {transcribed_text.strip()}")
                print(f"Emotion: {emotion} | Fluency: {fluency} | Pronunciation: {pronunciation} | Silvero:{silvero}")
                chunk_index += 1

    except WebSocketDisconnect:
        logging.info(f"[WS] {username} forcibly disconnected.")

    finally:
        await loop.run_in_executor(None, merge_chunks, chunk_files, final_output)
        logging.info(f"[Output] Final audio saved: {final_output}")

        with open(transcript_path, "w", encoding="utf-8") as f:
            f.write(" ".join(text_output).strip())
        logging.info(f"[Output] Transcript saved: {transcript_path}")

        for file in chunk_files:
            try:
                os.remove(file)
            except Exception as e:
                logging.info(f"[Warning] Failed to remove {file}: {e}")

def merge_chunks(chunk_files, final_output):
    logging.info("[Merge] Merging audio chunks...")
    combined = AudioSegment.empty()
    for file in chunk_files:
        audio = AudioSegment.from_file(file, format="wav")
        combined += audio
    combined.export(final_output, format="wav")
    logging.info("[Merge] Merged audio file saved.")


@app.get("/get-tts-audio")
def get_tts_audio(username: str):
    folder = os.path.join("text_to_speech_audio_folder", username)
    file_path = os.path.join(folder, f"{username}_output.wav")

    timeout = 60
    poll_interval = 2
    waited = 0

    while waited < timeout:
        if os.path.exists(file_path):
            return FileResponse(file_path, media_type="audio/wav", filename=f"{username}_output.wav")
        time.sleep(poll_interval)
        waited += poll_interval

    raise HTTPException(status_code=408, detail="Audio file not generated within 1 minute.")



genai.configure(api_key=GOOGLE_API_KEY)
gemini_model = genai.GenerativeModel("gemini-1.5-pro")

pc = Pinecone(api_key=PINECONE_API_KEY)

if PINECONE_INDEX_NAME not in pc.list_indexes().names():
    pc.create_index(
        name=PINECONE_INDEX_NAME,
        dimension=1024,
        metric="cosine",
        pods=1,
        pod_type="p1.x1"
    )
    logging.info(f"Created new Pinecone index {PINECONE_INDEX_NAME} with dimension 1024")

index = pc.Index(PINECONE_INDEX_NAME)

embedding_model = HuggingFaceEmbeddings(
    model_name=EMBEDDING_MODEL_NAME,
    model_kwargs={'device': 'cuda' if torch.cuda.is_available() else 'cpu'}
)


SUPPORTED_IMAGE_FORMATS = [".png", ".jpg", ".jpeg"]
SUPPORTED_TEXT_FORMATS = [".txt"]
SUPPORTED_PDF_FORMATS = [".pdf"]
SUPPORTED_DOC_FORMATS = [".docx"]
SUPPORTED_PPT_FORMATS = [".pptx"]
SUPPORTED_XLS_FORMATS = [".xlsx"]


def render_text_to_image(text: str, width=800, font_size=18) -> Image.Image:
    font = ImageFont.load_default()
    lines = []
    dummy_img = Image.new("RGB", (width, 1000))
    draw = ImageDraw.Draw(dummy_img)

    words = text.split()
    line = ""
    for word in words:
        test_line = f"{line} {word}".strip()
        bbox = draw.textbbox((0, 0), test_line, font=font)
        w = bbox[2] - bbox[0]
        if w < width - 40:
            line = test_line
        else:
            lines.append(line)
            line = word
    lines.append(line)

    height = font_size * len(lines) + 50
    img = Image.new("RGB", (width, height), color="white")
    draw = ImageDraw.Draw(img)
    y = 20
    for line in lines:
        draw.text((20, y), line, font=font, fill="black")
        y += font_size
    return img


def file_to_images(file_path: str) -> List[Image.Image]:
    ext = os.path.splitext(file_path)[1].lower()
    images = []

    if ext in SUPPORTED_PDF_FORMATS:
        images = convert_from_path(file_path, dpi=200)

    elif ext in SUPPORTED_IMAGE_FORMATS:
        images = [Image.open(file_path)]

    elif ext in SUPPORTED_DOC_FORMATS:
        try:
            text = docx2txt.process(file_path)
            img = render_text_to_image(text)
            images = [img]
        except Exception as e:
            raise HTTPException(status_code=422, detail=f"DOCX to image failed: {e}")

    elif ext in SUPPORTED_PPT_FORMATS:
        prs = Presentation(file_path)
        for i, slide in enumerate(prs.slides):
            text = "\n".join(shape.text for shape in slide.shapes if hasattr(shape, "text"))
            img = render_text_to_image(text)
            images.append(img)

    elif ext in SUPPORTED_XLS_FORMATS:
        try:
            excel = pd.read_excel(file_path, sheet_name=None)
            for sheet_name, df in excel.items():
                text = f"Sheet: {sheet_name}\n{df.to_string(index=False)}"
                img = render_text_to_image(text)
                images.append(img)
        except Exception as e:
            raise HTTPException(status_code=422, detail=f"XLSX to image failed: {e}")

    elif ext in SUPPORTED_TEXT_FORMATS:
        with open(file_path, 'r', encoding='utf-8') as f:
            content = f.read()
        img = render_text_to_image(content)
        images = [img]

    else:
        raise HTTPException(status_code=400, detail="Unsupported file format")

    return images


def extract_text_from_any_file(file_path: str) -> str:
    images = file_to_images(file_path)
    all_text = ""

    for idx, image in enumerate(images):
        try:
            response = gemini_model.generate_content([
                "Extract all the text from this image accurately, including tables and special formatting.",
                image
            ])
            text = response.text.strip()
            print(text)
            all_text += f"\n\n--- Page/Image {idx + 1} ---\n{text}"
        except Exception as e:
            logging.error(f"OCR failed on image {idx + 1}: {e}")
            all_text += f"\n\n--- Page/Image {idx + 1} FAILED ---"

    return all_text



@app.post("/upload/")
async def upload_file(
    file: UploadFile = File(...),
    student_class: str = Form(...),
    subject: str = Form(...),
    curriculum: str = Form(...)
):
    try:
        folder = f"uploads/{curriculum}/{student_class}/{subject}"
        os.makedirs(folder, exist_ok=True)
        file_path = os.path.join(folder, file.filename)

        with open(file_path, "wb") as buffer:
            shutil.copyfileobj(file.file, buffer)

        extracted_text = extract_text_from_any_file(file_path)

        if not extracted_text.strip():
            raise HTTPException(status_code=422, detail="No text extracted.")

        text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=1000,
            chunk_overlap=200,
            length_function=len
        )
        chunks = text_splitter.split_text(extracted_text)

        metadatas = [
            {
                "curriculum": curriculum,
                "student_class": student_class,
                "subject": subject,
                "filename": file.filename,
                "type": "ocr_file",
                "id": str(uuid.uuid4()),
                "text": chunk
            } for chunk in chunks
        ]

        vectorstore = PineconeVectorStore.from_existing_index(
            index_name=PINECONE_INDEX_NAME,
            embedding=embedding_model,
            text_key="text"
        )

        vectorstore.add_texts(texts=chunks, metadatas=metadatas)

        stats = index.describe_index_stats()
        serializable_stats = {
            "dimension": stats["dimension"],
            "index_fullness": stats["index_fullness"],
            "namespaces": {
                ns: {"vector_count": data["vector_count"]}
                for ns, data in stats["namespaces"].items()
            },
            "total_vector_count": stats["total_vector_count"]
        }

        logging.info(f"Upload successful. Total vectors: {stats['total_vector_count']}")
        return {
            "status": "success",
            "filepath": file_path,
            "extracted_text": extracted_text[:500] + "..."
        }

    except HTTPException as e:
        raise e

    except Exception as e:
        logging.error(f"Upload failed: {e}", exc_info=True)
        raise HTTPException(status_code=500, detail="Internal Server Error")


model_name = OllamaLLM(model="mistral")


@app.post("/chat/")
async def chat(request: ChatRequest):
    try:
        question = request.question.strip()
        subject = request.subject.strip()
        curriculum = request.curriculum.strip()

        vectorstore = PineconeVectorStore.from_existing_index(
            index_name=PINECONE_INDEX_NAME,
            embedding=embedding_model,
            text_key="text"
        )

        retriever = vectorstore.as_retriever(search_kwargs={
            "k": 10,
            "filter": {
                "subject": subject,
                "curriculum": curriculum
            }
        })

        prompt_template = PromptTemplate.from_template("""
        You are an expert educator providing clear, concise answers to students.
        Extract the most relevant information to answer the question using ONLY the provided context.

        Follow these rules:
        1. Answer in complete, well-structured sentences.
        2. Do not mention page numbers or document structure.
        3. If context doesn't contain any content, say "This information is not in our materials."
        4. Be factual and avoid speculation.
        5. Use proper grammar and spelling.
        6. Keep your answer concise and to the point.
        7. Do not include '\\n' or '*' in your output.
        8. Do not include escape characters like \\n, \\, \", \n or any slashes.
        9. Do not use markdown symbols like '*', '-', '`', or backslashes.

        Context: {context}
        Question: {question}
        Answer:
        """)

        qa_chain = RetrievalQA.from_chain_type(
            llm=model_name,
            chain_type="stuff",
            retriever=retriever,
            chain_type_kwargs={"prompt": prompt_template}
        )

        result = qa_chain.run(question)

        return {
            "question": question,
            "answer": result
        }

    except Exception as e:
        logging.error(f"Chat error: {str(e)}", exc_info=True)
        return {
            "question": request.question,
            "answer": "An error occurred while processing your request.",
            "error": str(e)
        }

@app.get("/test")
def welcome_page():
    return {"Message": "Welcome the ai speech module page."}

    


